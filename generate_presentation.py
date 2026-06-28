"""
Generate PowerPoint Presentation - Clean Light Theme
Best practices: light bg, large fonts, minimal text, one idea per slide.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pathlib import Path

FIGURES = Path("results/figures")

# Light professional color scheme
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BG_WHITE = RGBColor(0xFA, 0xFA, 0xFA)
DARK = RGBColor(0x1A, 0x23, 0x7E)
BLUE = RGBColor(0x15, 0x65, 0xC0)
GRAY = RGBColor(0x61, 0x61, 0x61)
LIGHT_GRAY = RGBColor(0x9E, 0x9E, 0x9E)
GREEN = RGBColor(0x2E, 0x7D, 0x32)
RED = RGBColor(0xC6, 0x28, 0x28)


def bg(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG_WHITE


def txt(slide, l, t, w, h, text, size=24, bold=False, color=DARK, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = "Calibri"
    p.alignment = align
    return tf


def bullets(slide, items, l=0.8, t=2.2, size=22, color=GRAY):
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(8.4), Inches(4.5))
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = Pt(12)


def title_bar(slide, title):
    bg(slide)
    # Blue bar at top
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  Inches(0), Inches(0), Inches(10), Inches(1.1))
    bar.fill.solid()
    bar.fill.fore_color.rgb = DARK
    bar.line.fill.background()
    txt(slide, 0.6, 0.2, 8.8, 0.7, title, size=32, bold=True, color=WHITE)


def img(slide, path, l, t, w, h=None):
    p = FIGURES / path
    if p.exists():
        if h:
            slide.shapes.add_picture(str(p), Inches(l), Inches(t), Inches(w), Inches(h))
        else:
            slide.shapes.add_picture(str(p), Inches(l), Inches(t), width=Inches(w))


def build():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    B = prs.slide_layouts[6]

    # ===== SLIDE 1: TITLE =====
    s = prs.slides.add_slide(B)
    bg(s)
    img(s, "cognitive_squad_banner_v2.png", 0.5, 0.3, 9, 2.5)
    txt(s, 0.5, 3.0, 9, 1.0,
        "EEG-Based Emotion Recognition Using\nMachine Learning and Deep Learning",
        size=28, bold=True, color=DARK, align=PP_ALIGN.CENTER)
    txt(s, 0.5, 4.5, 9, 0.5,
        "D.S.M. Perera  |  E.A.R. Fonseka  |  J.A.L. Manduli  |  K.M.H. Bandara",
        size=14, color=GRAY, align=PP_ALIGN.CENTER)
    txt(s, 0.5, 5.1, 9, 0.4,
        "MSc in Information Technology | Sri Lanka Institute of Information Technology",
        size=13, color=BLUE, align=PP_ALIGN.CENTER)
    txt(s, 0.5, 5.6, 9, 0.4,
        "Subject: Artificial Intelligence | Group Project",
        size=12, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

    # ===== SLIDE 2: PROBLEM =====
    s = prs.slides.add_slide(B)
    title_bar(s, "Why Emotion Recognition from EEG?")
    bullets(s, [
        "Emotions drive decision-making, learning, and mental health",
        "EEG captures brain activity directly and non-invasively",
        "Applications: mental health, adaptive learning, HCI",
        "Gap: No systematic ML vs DL comparison exists",
    ], t=1.8, size=24)

    # ===== SLIDE 3: OBJECTIVE =====
    s = prs.slides.add_slide(B)
    title_bar(s, "Research Objective")
    txt(s, 1.0, 2.5, 8, 2.0,
        "Compare 6 classification models\n(3 Traditional ML + 3 Deep Learning)\nfor EEG emotion recognition\nunder identical experimental conditions",
        size=28, color=DARK, align=PP_ALIGN.CENTER)
    txt(s, 1.0, 5.0, 8, 0.5,
        "Target: Classify Positive, Neutral, and Negative emotions from brainwaves",
        size=18, color=BLUE, align=PP_ALIGN.CENTER)

    # ===== SLIDE 4: DATASET =====
    s = prs.slides.add_slide(B)
    title_bar(s, "Dataset")
    bullets(s, [
        "Muse EEG Headband (4 channels)",
        "2 subjects | 3 emotions | 2,132 samples",
        "2,548 pre-extracted features per sample",
        "Balanced classes (~33% each)",
    ], t=1.5, size=22)
    img(s, "class_distribution.png", 5.5, 3.0, 4.2, 3.2)

    # ===== SLIDE 5: ARCHITECTURE =====
    s = prs.slides.add_slide(B)
    title_bar(s, "System Architecture")
    img(s, "system_architecture.png", 0.3, 1.5, 9.4, 3.5)
    txt(s, 0.5, 5.3, 9, 0.5,
        "End-to-end pipeline: Data > Preprocessing > 6 Models > Evaluation > Comparison",
        size=16, color=GRAY, align=PP_ALIGN.CENTER)

    # ===== SLIDE 6: MODELS =====
    s = prs.slides.add_slide(B)
    title_bar(s, "6 Models Evaluated")
    # Left: ML
    txt(s, 0.8, 1.5, 4, 0.5, "Traditional ML", size=22, bold=True, color=GREEN)
    bullets(s, [
        "SVM (RBF Kernel)",
        "Random Forest (300 trees)",
        "KNN (k=7)",
    ], l=0.8, t=2.2, size=20)
    # Right: DL
    txt(s, 5.5, 1.5, 4, 0.5, "Deep Learning", size=22, bold=True, color=BLUE)
    box = s.shapes.add_textbox(Inches(5.5), Inches(2.2), Inches(4.2), Inches(3))
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(["MLP (1.47M params)", "CNN-1D (44K params)", "LSTM (580K params)"]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(20)
        p.font.color.rgb = GRAY
        p.font.name = "Calibri"
        p.space_after = Pt(12)

    # ===== SLIDE 7: RESULTS =====
    s = prs.slides.add_slide(B)
    title_bar(s, "Results: All Models > 93% Accuracy")
    img(s, "model_comparison_accuracy.png", 0.5, 1.5, 9, 5.2)

    # ===== SLIDE 8: BEST RESULTS =====
    s = prs.slides.add_slide(B)
    title_bar(s, "Best Accuracy: 98.83%")
    txt(s, 1.0, 2.0, 8, 1.0,
        "CNN-1D  =  Random Forest  =  98.83%",
        size=36, bold=True, color=DARK, align=PP_ALIGN.CENTER)
    txt(s, 1.0, 3.5, 8, 0.5,
        "MLP: 98.36%   |   LSTM: 98.36%   |   SVM: 97.66%   |   KNN: 93.44%",
        size=20, color=GRAY, align=PP_ALIGN.CENTER)
    txt(s, 1.0, 4.8, 8, 1.0,
        "Key Insight: Feature engineering quality matters\nmore than model complexity",
        size=22, color=BLUE, align=PP_ALIGN.CENTER)

    # ===== SLIDE 9: CROSS-VALIDATION =====
    s = prs.slides.add_slide(B)
    title_bar(s, "10-Fold Cross-Validation")
    img(s, "cv_boxplot.png", 0.5, 1.3, 5.5, 4.5)
    txt(s, 6.3, 2.0, 3.5, 0.5, "Random Forest", size=22, bold=True, color=DARK)
    txt(s, 6.3, 2.6, 3.5, 0.5, "98.69% (+/-0.81%)", size=28, bold=True, color=GREEN)
    txt(s, 6.3, 3.5, 3.5, 0.5, "SVM", size=22, bold=True, color=DARK)
    txt(s, 6.3, 4.1, 3.5, 0.5, "97.94% (+/-0.92%)", size=24, color=BLUE)
    txt(s, 6.3, 5.0, 3.5, 0.5, "KNN", size=22, bold=True, color=DARK)
    txt(s, 6.3, 5.6, 3.5, 0.5, "94.18% (+/-1.87%)", size=24, color=GRAY)

    return prs


def add_analysis(prs):
    B = prs.slide_layouts[6]

    # ===== SLIDE 10: CONFUSION MATRICES =====
    s = prs.slides.add_slide(B)
    title_bar(s, "Per-Class Performance")
    img(s, "all_confusion_matrices.png", 0.3, 1.3, 9.4, 4.0)
    txt(s, 0.5, 5.5, 9, 0.5,
        "Neutral: easiest to classify  |  Positive vs Negative: main confusion source",
        size=16, color=GRAY, align=PP_ALIGN.CENTER)

    # ===== SLIDE 11: FEATURE IMPORTANCE =====
    s = prs.slides.add_slide(B)
    title_bar(s, "Most Important EEG Features")
    img(s, "feature_importance_top30.png", 0.3, 1.3, 9.4, 5.5)

    # ===== SLIDE 12: PCA =====
    s = prs.slides.add_slide(B)
    title_bar(s, "PCA: Clear Class Separation")
    img(s, "pca_2d_visualization.png", 1.5, 1.3, 7, 5.5)

    # ===== SLIDE 13: DISCUSSION =====
    s = prs.slides.add_slide(B)
    title_bar(s, "Key Findings")
    bullets(s, [
        "Random Forest = CNN-1D at 98.83%",
        "Feature quality > model complexity",
        "CNN-1D: 44K params beats MLP: 1.47M params",
        "Our results surpass published benchmarks",
    ], t=1.8, size=24, color=DARK)

    # ===== SLIDE 14: LIMITATIONS =====
    s = prs.slides.add_slide(B)
    title_bar(s, "Limitations")
    bullets(s, [
        "Only 2 subjects (limited generalizability)",
        "Pre-extracted features (no raw signal learning)",
        "4 EEG channels (limited spatial resolution)",
        "Controlled lab setting (not real-world)",
    ], t=1.8, size=24, color=GRAY)

    # ===== SLIDE 15: FUTURE WORK =====
    s = prs.slides.add_slide(B)
    title_bar(s, "Future Work")
    bullets(s, [
        "Larger datasets (DEAP, SEED)",
        "Raw EEG signal processing (end-to-end)",
        "Real-time classification system",
        "Transformer architectures",
        "Cross-subject transfer learning",
    ], t=1.8, size=24, color=DARK)

    # ===== SLIDE 16: CONCLUSION =====
    s = prs.slides.add_slide(B)
    title_bar(s, "Conclusion")
    txt(s, 0.8, 2.0, 8.4, 3.0,
        "98.83% accuracy achieved for\n3-class EEG emotion recognition\n\n"
        "Both traditional ML and deep learning\nare effective with quality features\n\n"
        "Consumer-grade EEG devices are viable\nfor emotion recognition applications",
        size=26, color=DARK, align=PP_ALIGN.CENTER)

    # ===== SLIDE 17: Q&A =====
    s = prs.slides.add_slide(B)
    bg(s)
    img(s, "cognitive_squad_banner_v2.png", 0.5, 0.5, 9, 2.5)
    txt(s, 0.5, 3.5, 9, 1.0,
        "Questions?",
        size=44, bold=True, color=DARK, align=PP_ALIGN.CENTER)
    txt(s, 0.5, 4.8, 9, 0.5,
        "Thank you for your attention!",
        size=20, color=BLUE, align=PP_ALIGN.CENTER)
    txt(s, 0.5, 5.6, 9, 0.5,
        "Cognitive Squad | MSc in Information Technology | SLIIT",
        size=14, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

    return prs


def main():
    print("Generating presentation (light theme)...")
    prs = build()
    prs = add_analysis(prs)
    out = "paper/Cognitive_Squad_Presentation_v3.pptx"
    prs.save(out)
    print(f"Saved: {out}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
