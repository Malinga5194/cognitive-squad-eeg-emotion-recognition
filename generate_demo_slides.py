"""
Generate Demo Video Slides - D.S.M. Perera (Individual Submission)
Two slides: Title + Closing.
Formal light theme with Cognitive Squad logo and accent colors.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pathlib import Path

FIGURES = Path("results/figures")
LOGO = FIGURES / "cognitive_squad_logo_v2.png"

# Formal color palette
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BG_WHITE = RGBColor(0xFB, 0xFB, 0xFD)
NAVY = RGBColor(0x14, 0x1E, 0x61)      # dark navy header
DARK = RGBColor(0x1A, 0x23, 0x7E)
BLUE = RGBColor(0x15, 0x65, 0xC0)      # accent blue
GRAY = RGBColor(0x55, 0x55, 0x55)
LIGHT_GRAY = RGBColor(0x90, 0x90, 0x90)
GREEN = RGBColor(0x2E, 0x7D, 0x32)


def bg(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG_WHITE


def txt(slide, l, t, w, h, text, size=24, bold=False, italic=False,
        color=DARK, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(size)
        p.font.bold = bold
        p.font.italic = italic
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.alignment = align
    return tf


def rect(slide, l, t, w, h, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Inches(l), Inches(t), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def add_logo(slide, l, t, w):
    if LOGO.exists():
        slide.shapes.add_picture(str(LOGO), Inches(l), Inches(t), width=Inches(w))


def build():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    B = prs.slide_layouts[6]

    # ============ SLIDE 1: TITLE ============
    s = prs.slides.add_slide(B)
    bg(s)
    # Navy header band
    rect(s, 0, 0, 10, 1.7, NAVY)
    # Blue accent line under header
    rect(s, 0, 1.7, 10, 0.08, BLUE)
    # Logo on header
    add_logo(s, 0.45, 0.3, 1.1)
    txt(s, 1.7, 0.45, 7.8, 0.6, "COGNITIVE SQUAD",
        size=30, bold=True, color=WHITE)
    txt(s, 1.72, 1.12, 7.8, 0.4, "EEG-Based Emotion Recognition System",
        size=15, italic=True, color=RGBColor(0xC5, 0xCD, 0xF0))

    # Main title
    txt(s, 0.7, 2.4, 8.6, 1.4,
        "EEG-Based Emotion Recognition Using\nMachine Learning and Deep Learning",
        size=30, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

    # Presenter block
    rect(s, 2.7, 4.15, 4.6, 0.05, BLUE)
    txt(s, 0.7, 4.35, 8.6, 0.5, "D.S.M. Perera",
        size=24, bold=True, color=DARK, align=PP_ALIGN.CENTER)
    txt(s, 0.7, 4.9, 8.6, 0.4, "Team Lead / System Architect  |  MS26906294",
        size=15, color=BLUE, align=PP_ALIGN.CENTER)

    # Footer details
    txt(s, 0.7, 5.7, 8.6, 0.4,
        "MSc in Information Technology  |  Sri Lanka Institute of Information Technology",
        size=14, color=GRAY, align=PP_ALIGN.CENTER)
    txt(s, 0.7, 6.15, 8.6, 0.4,
        "Subject: Artificial Intelligence  |  Group Project (Individual Demonstration)",
        size=13, italic=True, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

    # ============ SLIDE 2: CLOSING ============
    s = prs.slides.add_slide(B)
    bg(s)
    rect(s, 0, 0, 10, 1.7, NAVY)
    rect(s, 0, 1.7, 10, 0.08, BLUE)
    add_logo(s, 0.45, 0.3, 1.1)
    txt(s, 1.7, 0.55, 7.8, 0.7, "Thank You", size=34, bold=True, color=WHITE)

    # Key achievement highlight
    txt(s, 0.7, 2.3, 8.6, 0.6, "98.83% Accuracy",
        size=40, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
    txt(s, 0.7, 3.15, 8.6, 0.5,
        "3-Class Emotion Recognition from Brainwave Signals",
        size=18, color=GRAY, align=PP_ALIGN.CENTER)

    # Contribution summary
    rect(s, 2.7, 3.95, 4.6, 0.05, BLUE)
    txt(s, 0.7, 4.15, 8.6, 1.1,
        "System Architecture  •  Transformer, CNN-LSTM & LSTM Models  •  Web Dashboard",
        size=16, color=DARK, align=PP_ALIGN.CENTER)

    txt(s, 0.7, 5.4, 8.6, 0.4, "D.S.M. Perera  |  Team Lead / System Architect",
        size=15, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    txt(s, 0.7, 5.9, 8.6, 0.4,
        "Cognitive Squad  |  MSc in Information Technology  |  SLIIT",
        size=13, color=GRAY, align=PP_ALIGN.CENTER)
    txt(s, 0.7, 6.35, 8.6, 0.4,
        "Subject: Artificial Intelligence",
        size=12, italic=True, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

    return prs


def main():
    out_candidates = [
        "paper/demo_slides_Perera_v4.pptx",
        "paper/demo_slides_Perera_v5.pptx",
        "paper/demo_slides_Perera_v6.pptx",
    ]
    prs = build()
    last_err = None
    for out in out_candidates:
        try:
            prs.save(out)
            print(f"Saved: {out}")
            print(f"Slides: {len(prs.slides)}")
            return
        except PermissionError as e:
            last_err = e
            print(f"Locked (open in PowerPoint): {out} - trying next name")
    raise last_err


if __name__ == "__main__":
    main()
